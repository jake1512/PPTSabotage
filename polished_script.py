import os
import pptx
import random
from tkinter import *
from tkinter import filedialog as fd

TEXT_SHAPE_TYPES = ['PLACEHOLDER', 'TEXT_BOX', 'AUTO_SHAPE']
extensions = ['*.ppt', '*.pptx']
fonts = [
        'Marlett'
        ,'MS Outlook'
        ,'MS Reference Specialty'
        ,'Wingdings'
        ,'Webdings'
        ,'Symbol'
        ,'MT Extra'
        ,''
        ]

def shuffle_word(word):
    word = list(word)
    for i in range(len(word)):
        word[i] = chr(ord(word[i]) + 1)
    random.shuffle(word)
    return ''.join(word)

def random_font():
    font = random.choice(fonts)
    return font

def get_slide_count(prs):
    slide_count = 0
    for slide in prs.slides:
        slide_count += 1
    return slide_count

def get_random_slides(slide_count):
    percentage = 1 / 5
    affected_count = round((slide_count - 2) * percentage)
    affected_slides = random.sample(range(3, slide_count), affected_count)
    affected_slides = sorted(affected_slides)
    return affected_slides

def text_manipulation(text_frame, basetext):
    newtext = shuffle_word(basetext)
    for i in range(len(text_frame.paragraphs)):
        p = text_frame.paragraphs[i]
        p.clear()
        run = p.add_run()
        run.text = newtext
        newfont = random_font()
        if(newfont != ''):
            font = run.font
            font.name = newfont

def edit(path):
    prs = pptx.Presentation(path)
    victim_array = get_random_slides(get_slide_count(prs))
    print('Affected slides:', victim_array)
    slides = prs.slides
    for idx in victim_array:
        victim = slides[idx]
        for shape in victim.shapes:
            for shapetype in TEXT_SHAPE_TYPES:

                if shapetype in str(shape.shape_type):
                    text_frame = shape.text_frame
                    basetext = shape.text
                    text_manipulation(text_frame, basetext)

                if 'TABLE' in str(shape.shape_type):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text_frame = cell.text_frame
                            basetext = cell.text
                            text_manipulation(text_frame, basetext)
                
                if 'PICTURE' in str(shape.shape_type):
                    shape.rotation = 180     
            
    head_tail = os.path.split(path)
    targetPath = head_tail[0] + '/' + 'fixed'
    os.makedirs(targetPath, exist_ok=False)
    newfile = targetPath + '/' + head_tail[1][:-5] + '.' + head_tail[1][-5:]
    print(newfile)
    prs.save(newfile)

def convert_to_float(frac_str):
    try:
        return float(frac_str)
    except ValueError:
        num, denom = frac_str.split('/')
        try:
            leading, num = num.split(' ')
            whole = float(leading)
        except ValueError:
            whole = 0
        frac = float(num) / float(denom)
        return whole - frac if whole < 0 else whole + frac

def main():
    win = Tk()
    win.geometry("800x350")
    
    files = fd.askopenfilenames(parent=win, title='Choose a File', filetypes=[('PowerPoint', extensions)])
    
    has_files_text = StringVar()
    if files != '':
        has_files_text.set('Chosen files: \n')
        has_file = Label(win, textvariable=has_files_text)
        has_file.pack()      

        for filepath in files:
            filename_text = StringVar()
            filename_text.set(filepath)
            filename = Label(win, textvariable=filename_text)
            filename.pack()
            # print('Editing ' + filepath)
            # edit(filepath)
            print('\n')
    else:
        has_files_text.set('No files chosen')
        has_file = Label(win, textvariable=has_files_text)
        has_file.pack()
    
    win.mainloop()
    input("Press enter to exit")

if __name__ == '__main__':
    main()