import pptx
prs = pptx.Presentation('E:/Work/Programming_Stuffs/Python/PPTSabotage/SML Einstein HCP Profile _ Call script .pptx')
slides = prs.slides
victim = slides[16]

for shape in victim.shapes:
    print('%s, %s' % (shape.shape_type, shape.name))

# for shape in victim.placeholders:
#     print('%s' % (shape.name))