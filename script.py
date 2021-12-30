import pptx
from random import shuffle


def shuffle_word(word):
    word = list(word)
    for i in range(len(word)):
        word[i] = chr(ord(word[i]) + 1)
    shuffle(word)
    return ''.join(word)


prs = pptx.Presentation('E:/Work/Programming_Stuffs/Python/PPTSabotage/SML Einstein HCP Profile _ Call script .pptx')
slides = prs.slides
victim = slides[17]
# for victim in slides:
for shape in victim.shapes:

    # if shape.is_placeholder:
    #     text_frame = shape.text_frame
    #     for i in range(len(text_frame.paragraphs)):
    #         p = text_frame.paragraphs[i]
    #         basetext = shape.text
    #         newtext = shuffle_word(basetext)
    #         p.clear()
    #         run = p.add_run()
    #         run.text = newtext
    #         # font = run.font
    #         # font.name = "Times New Roman"

    # if 'PICTURE' in str(shape.shape_type):
    #     shape.rotation = 180

    # if 'TEXT_BOX' in str(shape.shape_type):
    #     text_frame = shape.text_frame
    #     basetext = shape.text
    #     newtext = shuffle_word(basetext)
    #     for i in range(len(text_frame.paragraphs)):
    #         p = text_frame.paragraphs[i]
    #         p.clear()
    #         run = p.add_run()
    #         run.text = newtext
    #         # font = run.font
    #         # font.name = "Times New Roman"

    # if 'AUTO_SHAPE' in str(shape.shape_type):
    #     text_frame = shape.text_frame
    #     # print(text_frame.paragraphs)
    #     basetext = shape.text
    #     newtext = shuffle_word(basetext)
        # for i in range(len(text_frame.paragraphs)):
        #     p = text_frame.paragraphs[i]
        #     # print(basetext)
        #     p.clear()
        #     run = p.add_run()
        #     run.text = newtext
        #     # font = run.font
        #     # font.name = "Times New Roman"
    
    if 'TABLE' in str(shape.shape_type):
        for row in shape.table.rows:
            for cell in row.cells:
                # print(cell.text_frame.paragraphs)
                text_frame = cell.text_frame
                basetext = cell.text
                newtext = shuffle_word(basetext)
                for i in range(len(text_frame.paragraphs)):
                    p = text_frame.paragraphs[i]
                    p.clear()
                    run = p.add_run()
                    run.text = newtext
                    # font = run.font
                    # font.name = "Times New Roman"
                
        
prs.save('./PPTSabotage/SML Einstein HCP Profile _ Call script_fkd_up.pptx')







